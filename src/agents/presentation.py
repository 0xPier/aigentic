"""Presentation Agent for creating PowerPoint and PDF presentations."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from typing import Dict, Any, List
from datetime import datetime
import json

from src.agents.base import LLMAgent, AgentContext, AgentResult


class PresentationAgent(LLMAgent):
    """Agent specialized in creating presentations and slide decks."""
    
    def __init__(self):
        super().__init__(
            name="presentation",
            description="Creates PowerPoint presentations, slide decks, and PDF documents",
            capabilities=[
                "powerpoint_creation",
                "slide_design",
                "pdf_generation",
                "presentation_structure",
                "visual_storytelling"
            ]
        )
    
    def get_required_integrations(self) -> List[str]:
        """Presentation agent doesn't require specific integrations."""
        return []
    
    async def execute(self, context: AgentContext) -> AgentResult:
        """Execute presentation creation task."""
        start_time = datetime.utcnow()
        
        try:
            # Parse presentation requirements
            presentation_spec = await self._parse_presentation_requirements(context.query)
            
            # Create presentation structure
            structure = await self._create_presentation_structure(presentation_spec, context)
            
            # Generate presentation content
            content = await self._generate_presentation_content(structure, context)
            
            # Create PowerPoint presentation
            pptx_file = await self._create_powerpoint(content, presentation_spec)
            
            execution_time = (datetime.utcnow() - start_time).total_seconds()
            
            result = AgentResult(
                success=True,
                data={
                    "presentation_spec": presentation_spec,
                    "structure": structure,
                    "content": content,
                    "slide_count": len(content.get("slides", [])),
                    "theme": presentation_spec.get("theme", "professional")
                },
                message="Presentation created successfully",
                execution_time=execution_time,
                output_files=[pptx_file] if pptx_file else []
            )
            
            # Save to memory
            self.save_memory(
                memory_type="presentation",
                content={
                    "topic": presentation_spec.get("topic", ""),
                    "slide_count": len(content.get("slides", [])),
                    "theme": presentation_spec.get("theme", "professional")
                },
                context_tags=["presentation", "slides"],
                relevance_score=0.8
            )
            
            self.log_execution(context, result)
            return result
            
        except Exception as e:
            execution_time = (datetime.utcnow() - start_time).total_seconds()
            
            result = AgentResult(
                success=False,
                data={},
                message=f"Presentation creation failed: {str(e)}",
                execution_time=execution_time,
                error=str(e)
            )
            
            self.log_execution(context, result)
            return result
    
    async def _parse_presentation_requirements(self, query: str) -> Dict[str, Any]:
        """Parse presentation requirements from query."""
        system_prompt = """
        You are a presentation expert. Parse the requirements and determine:
        1. Presentation topic and purpose
        2. Target audience
        3. Presentation length (number of slides)
        4. Theme and style (professional, creative, minimal, etc.)
        5. Key sections to include
        6. Visual elements needed
        
        Respond in JSON format:
        {
            "topic": "presentation topic",
            "purpose": "inform|persuade|educate|pitch",
            "audience": "target audience",
            "slide_count": 10,
            "theme": "professional",
            "sections": ["introduction", "main_content", "conclusion"],
            "visual_elements": ["charts", "images", "bullet_points"]
        }
        """
        
        prompt = f"Parse presentation requirements for: {query}"
        
        try:
            response = await self.call_llm(prompt, system_prompt, temperature=0.3)
            return json.loads(response)
        except Exception as e:
            self.logger.error(f"Presentation parsing failed: {e}")
            return {
                "topic": query,
                "purpose": "inform",
                "audience": "general",
                "slide_count": 8,
                "theme": "professional",
                "sections": ["introduction", "main_content", "conclusion"],
                "visual_elements": ["bullet_points", "charts"]
            }
    
    async def _create_presentation_structure(self, spec: Dict[str, Any], context: AgentContext) -> Dict[str, Any]:
        """Create the overall structure of the presentation."""
        system_prompt = """
        You are a presentation structure expert. Create a detailed outline including:
        1. Slide-by-slide breakdown
        2. Content flow and narrative
        3. Visual elements for each slide
        4. Transitions and connections
        
        Respond in JSON format:
        {
            "outline": [
                {
                    "slide_number": 1,
                    "title": "slide title",
                    "type": "title|content|section|conclusion",
                    "key_points": ["point1", "point2"],
                    "visual_elements": ["chart", "image"]
                }
            ],
            "narrative_flow": "description of story flow",
            "key_messages": ["message1", "message2"]
        }
        """
        
        prompt = f"""
        Presentation Specification:
        {json.dumps(spec, indent=2)}
        
        Context: {context.query}
        
        Create a detailed presentation structure.
        """
        
        try:
            response = await self.call_llm(prompt, system_prompt, temperature=0.3)
            return json.loads(response)
        except Exception as e:
            self.logger.error(f"Structure creation failed: {e}")
            return {
                "outline": [
                    {"slide_number": 1, "title": "Title Slide", "type": "title", "key_points": [spec.get("topic", "Presentation")], "visual_elements": []},
                    {"slide_number": 2, "title": "Overview", "type": "content", "key_points": ["Agenda", "Objectives"], "visual_elements": ["bullet_points"]},
                    {"slide_number": 3, "title": "Main Content", "type": "content", "key_points": ["Key information"], "visual_elements": ["bullet_points"]},
                    {"slide_number": 4, "title": "Conclusion", "type": "conclusion", "key_points": ["Summary", "Next steps"], "visual_elements": ["bullet_points"]}
                ],
                "narrative_flow": "Introduction → Main content → Conclusion",
                "key_messages": ["Main topic overview", "Key insights", "Action items"]
            }
    
    async def _generate_presentation_content(self, structure: Dict[str, Any], context: AgentContext) -> Dict[str, Any]:
        """Generate detailed content for each slide."""
        slides = []
        
        for slide_outline in structure.get("outline", []):
            slide_content = await self._generate_slide_content(slide_outline, context)
            slides.append(slide_content)
        
        return {
            "slides": slides,
            "theme_colors": ["#1f4e79", "#ffffff", "#d9e2ec"],  # Professional blue theme
            "font_family": "Calibri"
        }
    
    async def _generate_slide_content(self, outline: Dict[str, Any], context: AgentContext) -> Dict[str, Any]:
        """Generate content for a single slide."""
        system_prompt = f"""
        You are a presentation content writer. Create detailed content for a slide with:
        - Title: {outline.get('title', '')}
        - Type: {outline.get('type', 'content')}
        - Key points: {outline.get('key_points', [])}
        
        Provide:
        1. Refined slide title
        2. Bullet points or content text
        3. Speaker notes
        4. Visual suggestions
        
        Respond in JSON format:
        {{
            "title": "slide title",
            "content": ["bullet point 1", "bullet point 2"],
            "speaker_notes": "notes for presenter",
            "visual_suggestions": ["suggestion1", "suggestion2"]
        }}
        """
        
        prompt = f"""
        Slide outline: {json.dumps(outline, indent=2)}
        Overall context: {context.query}
        
        Generate detailed slide content.
        """
        
        try:
            response = await self.call_llm(prompt, system_prompt, temperature=0.4)
            slide_data = json.loads(response)
            slide_data["slide_number"] = outline.get("slide_number", 1)
            slide_data["type"] = outline.get("type", "content")
            return slide_data
        except Exception as e:
            self.logger.error(f"Slide content generation failed: {e}")
            return {
                "slide_number": outline.get("slide_number", 1),
                "title": outline.get("title", "Slide Title"),
                "type": outline.get("type", "content"),
                "content": outline.get("key_points", ["Content point"]),
                "speaker_notes": "Speaker notes for this slide",
                "visual_suggestions": ["Add relevant visual elements"]
            }
    
    async def _create_powerpoint(self, content: Dict[str, Any], spec: Dict[str, Any]) -> str:
        """Create PowerPoint presentation file."""
        try:
            # Create presentation
            prs = Presentation()
            
            # Set theme colors
            theme_colors = content.get("theme_colors", ["#1f4e79", "#ffffff", "#d9e2ec"])
            
            for slide_data in content.get("slides", []):
                slide_type = slide_data.get("type", "content")
                
                if slide_type == "title":
                    slide = self._create_title_slide(prs, slide_data, theme_colors)
                else:
                    slide = self._create_content_slide(prs, slide_data, theme_colors)
            
            # Save presentation
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"presentation_{timestamp}.pptx"
            filepath = os.path.join("uploads", filename)
            
            # Create uploads directory if it doesn't exist
            os.makedirs("uploads", exist_ok=True)
            
            prs.save(filepath)
            return filepath
            
        except Exception as e:
            self.logger.error(f"PowerPoint creation failed: {e}")
            return None
    
    def _create_title_slide(self, prs: Presentation, slide_data: Dict[str, Any], colors: List[str]) -> Any:
        """Create title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data.get("title", "Presentation Title")
        
        # Set subtitle if content exists
        if slide.shapes.placeholders[1] and slide_data.get("content"):
            subtitle = slide.shapes.placeholders[1]
            subtitle.text = "\n".join(slide_data.get("content", []))
        
        # Style the title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.color.rgb = RGBColor.from_string(colors[0].replace("#", ""))
        
        return slide
    
    def _create_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], colors: List[str]) -> Any:
        """Create content slide with bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data.get("title", "Slide Title")
        
        # Set content
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(slide_data.get("content", [])):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
        
        # Style the title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.color.rgb = RGBColor.from_string(colors[0].replace("#", ""))
        
        return slide
